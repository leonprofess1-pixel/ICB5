import os
import time
import random
import requests
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# [설정] 저장할 폴더 및 파일 이름
FOLDER_NAME = "명언_결과물"
IMAGE_FOLDER = "images"
FILE_NAME = "명언_출력용_최종완성.pptx"

# 폴더 생성
if not os.path.exists(FOLDER_NAME):
    os.makedirs(FOLDER_NAME)
if not os.path.exists(IMAGE_FOLDER):
    os.makedirs(IMAGE_FOLDER)

# 데이터 준비: (명언, 인물, 영어_검색키워드)
# * 영어 키워드를 사용해야 정확도가 높습니다.
quotes_data = [
    ("재능은 게임을 이기게 하지만, 팀워크와 이해관계는 진정한 우승을 가져다줍니다.", "마이클 조던", "basketball"),
    ("만남은 시작이고, 대화는 진보이며, 협력은 성공이라는 것을 항상 기억하십시오.", "헨리 포드", "handshake"),
    ("팀워크는 평범한 사람들이 비범한 결과를 낼 수 있도록 만드는 원동력입니다.", "앤드류 카네기", "teamwork"),
    ("우리가 들이마시는 공기처럼, 신뢰 없이는 어떤 팀도 생존할 수 없습니다.", "워렌 버핏", "trust"),
    ("당신의 전략이 아무리 뛰어나도, 혼자 게임을 한다면 언제나 팀에게 질 것입니다.", "리드 호프만", "chess"),
    ("훌륭한 팀은 서로 감추는 것이 없으며, 문제를 두려워하지 않고 함께 해결합니다.", "에드 캣멀", "meeting"),
    ("조직을 승리로 이끄는 힘은 25%의 실력과 75%의 팀워크에서 나옵니다.", "딕 버메일", "sports"),
    ("위대한 일은 결코 한 사람이 아니라, 헌신적인 팀에 의해 비로소 완성됩니다.", "스티브 잡스", "apple_computer"),
    ("개별적으로 우리는 한 방울의 물이지만, 함께 모이면 우리는 거대한 바다가 됩니다.", "류노스케", "ocean"),
    ("위험을 감수하지 않는 것이야말로 인생에서 가장 큰 위험이라는 것을 명심하십시오.", "스티븐 코비", "cliff"),
    ("어제와 똑같이 살면서 다른 미래를 기대하는 것은 정신병 초기 증세다.", "알베르트 아인슈타인", "clock"),
    ("배는 항구에 정박해 있을 때 가장 안전하지만, 그것이 배가 만들어진 존재 이유는 아닙니다.", "존 A. 셰드", "ship"),
    ("가장 강한 종이 살아남는 것이 아니라, 변화에 가장 잘 적응하는 종이 살아남습니다.", "찰스 다윈", "nature"),
    ("안락함의 지대에서 벗어나십시오. 편안함의 끝에서야 비로소 진정한 성장이 시작됩니다.", "로이 T. 베넷", "hiking"),
    ("아무것도 하지 않으면 의심과 공포가 생겨나고, 행동하면 자신감과 용기가 생겨납니다.", "데일 카네기", "courage"),
    ("미래를 예측하는 가장 훌륭하고 확실한 방법은 스스로 미래를 창조하는 것입니다.", "피터 드러커", "future"),
    ("실패는 다시 시작할 수 있는 기회이며, 이번에는 더 현명하게 시작할 수 있다는 뜻입니다.", "헨리 포드", "sunrise"),
    ("100퍼센트의 확신이 들 때까지 기다린다면, 그때는 이미 늦은 것입니다.", "하워드 슐츠", "running"),
    ("시도해보지 않고는 그 누구도 자신이 얼마만큼 해낼 수 있는지 알지 못합니다.", "푸블릴리우스 시루스", "challenge"),
    ("당신의 재능이 아니라 당신의 태도가 당신의 인생의 높이를 결정할 것입니다.", "지그 지글러", "stairs"),
    ("당신이 할 수 있다고 믿든 할 수 없다고 믿든, 당신이 믿는 대로 될 것입니다.", "헨리 포드", "belief"),
    ("비관론자는 모든 기회에서 어려움을 찾아내고, 낙관론자는 모든 어려움에서 기회를 찾아냅니다.", "윈스턴 처칠", "light"),
    ("행복해서 웃는 것이 아니라, 웃기 때문에 행복해진다는 사실을 잊지 마십시오.", "윌리엄 제임스", "smile"),
    ("얼굴을 언제나 태양으로 향하십시오. 그러면 그림자는 당신의 뒤로 물러날 것입니다.", "헬렌 켈러", "sunflower"),
    ("진정한 발견은 새로운 땅을 찾는 것이 아니라, 새로운 눈으로 세상을 보는 것입니다.", "마르셀 프루스트", "travel"),
    ("내가 헛되이 보낸 오늘은 어제 죽은 이가 그토록 갈망하던 내일임을 기억하십시오.", "소포클레스", "time"),
    ("위대한 업적은 대게 커다란 위험을 감수한 결과라는 것을 역사가 증명합니다.", "헤로도토스", "achievement"),
    ("생각하는 대로 살지 않으면, 머지않아 사는 대로 생각하게 될 것입니다.", "폴 발레리", "thinking"),
    ("성공한 사람이 되려 하기보다는 가치 있는 사람이 되기 위해 노력하십시오.", "알베르트 아인슈타인", "value"),
    ("성공이란 열정을 잃지 않고 실패에서 실패로 거듭 나아가는 능력입니다.", "윈스턴 처칠", "passion"),
    ("꿈을 꿀 수 있다면 이룰 수도 있습니다. 내 모든 것이 꿈에서 시작되었다는 것을 기억하십시오.", "월트 디즈니", "dream"),
    ("천재는 1퍼센트의 영감과 99퍼센트의 땀과 노력으로 만들어지는 것입니다.", "토마스 에디슨", "lightbulb"),
    ("우리는 우리가 반복적으로 하는 행동 그 자체입니다. 그러므로 탁월함은 행동이 아니라 습관입니다.", "아리스토텔레스", "habit"),
    ("계획 없는 목표는 한낱 꿈에 불과하며, 실행 없는 계획은 시간 낭비일 뿐입니다.", "생텍쥐페리", "map"),
    ("승리는 대개 끈기 있게 노력하는 자에게 돌아간다는 사실을 잊지 마십시오.", "빈스 롬바르디", "trophy"),
    ("장애물은 당신이 목표에서 눈을 뗐을 때 비로소 보이게 되는 무서운 것들입니다.", "헨리 포드", "obstacle"),
    ("누군가 그늘에 앉아 쉴 수 있는 건, 오래전 다른 누군가가 나무를 심었기 때문입니다.", "워렌 버핏", "tree"),
    ("인내는 쓰지만 그 열매는 달다는 것을 기억하고 끝까지 포기하지 마십시오.", "장 자크 루소", "fruit"),
    ("성공이 노력보다 먼저 나오는 곳은 이 세상에서 오직 영어 사전뿐입니다.", "비달 사순", "dictionary"),
    ("배움을 멈추는 사람은 스무 살이든 여든 살이든 늙은 것과 다름없습니다.", "헨리 포드", "reading"),
    ("리더십은 다음 세대를 위한 기회를 만들어내는 것이지, 다음 선거를 준비하는 것이 아닙니다.", "제임스 프리먼 클라크", "leadership"),
    ("남보다 더 잘하려고 하지 말고, 어제의 나보다 더 잘하려고 노력하십시오.", "윌리엄 포크너", "mirror"),
    ("질문을 잊지 않는 것이야말로 지식을 넓히고 지혜를 얻는 가장 빠른 길입니다.", "알베르트 아인슈타인", "question"),
    ("지식에 투자하는 것이야말로 언제나 가장 높은 이자를 지불해 준다는 것을 잊지 마십시오.", "벤자민 프랭클린", "investment"),
    ("시간을 지배할 줄 아는 사람은 인생을 지배할 줄 아는 사람입니다.", "에센바흐", "calendar"),
    ("가장 유능한 사람은 가장 배울 점이 많은 사람이라는 것을 잊지 마십시오.", "괴테", "student"),
    ("리더는 길을 아는 사람이고, 길을 가는 사람이며, 길을 보여주는 사람입니다.", "존 맥스웰", "compass"),
    ("아는 것만으로는 충분하지 않습니다. 적용해야 합니다. 의지만으로는 충분하지 않습니다. 실행해야 합니다.", "괴테", "action"),
    ("인생에서 가장 큰 영광은 넘어지지 않는 것이 아니라, 넘어질 때마다 다시 일어서는 데 있습니다.", "넬슨 만델라", "rodeo")
]

# --- [1] 이미지 다운로드 함수 (LoremFlickr 사용) ---
def download_image(keyword, index):
    filename = f"image_{index}.jpg"
    filepath = os.path.join(IMAGE_FOLDER, filename)
    
    # 이미 파일이 있으면 다시 받지 않음 (시간 절약)
    if os.path.exists(filepath):
        return filepath

    print(f"[{index+1}/{len(quotes_data)}] '{keyword}' 관련 이미지 다운로드 중...")
    
    # 키워드 기반 랜덤 이미지 URL (800x600 사이즈)
    # random 파라미터를 추가하여 매번 다른 이미지가 나오도록 함
    url = f"https://loremflickr.com/800/600/{keyword}?random={index}"
    
    try:
        response = requests.get(url, timeout=15)
        if response.status_code == 200:
            with open(filepath, 'wb') as f:
                f.write(response.content)
            return filepath
        else:
            print(f"  -> 다운로드 실패 (상태 코드: {response.status_code})")
            return None
    except Exception as e:
        print(f"  -> 오류 발생: {e}")
        return None

# --- [2] PPT 생성 함수 ---
def create_final_ppt():
    prs = Presentation()
    slide_width = Cm(39.5)
    slide_height = Cm(8)
    prs.slide_width = slide_width
    prs.slide_height = slide_height

    def get_random_pastel_color():
        return RGBColor(random.randint(220, 255), random.randint(220, 255), random.randint(220, 255))

    for i, (quote, author, keyword) in enumerate(quotes_data):
        # 1. 이미지 준비
        img_path = download_image(keyword, i)
        
        # 2. 슬라이드 추가
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # [배경] 랜덤 그라데이션
        fill = slide.background.fill
        fill.gradient()
        fill.gradient_stops[0].position = 0.0
        fill.gradient_stops[0].color.rgb = get_random_pastel_color()
        # 그라데이션 중지점이 부족할 경우 추가
        if len(fill.gradient_stops) < 2: 
            fill.gradient_stops.add_stop(1.0, get_random_pastel_color())
        else: 
            fill.gradient_stops[1].color.rgb = get_random_pastel_color()
        fill.gradient_angle = random.choice(range(0, 360, 45))

        # [테두리] 제단선
        border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_width, slide_height)
        border.fill.background() # 투명
        border.line.color.rgb = RGBColor(0, 0, 0)
        border.line.width = Pt(1)

        # --- [이미지 배치] ---
        # 요구사항: 좌측 배치, 가로 7cm, 세로 5cm
        img_w = Cm(7)
        img_h = Cm(5)
        img_left = Cm(1) # 왼쪽 여백 1cm
        img_top = (slide_height - img_h) / 2 # 수직 가운데 정렬

        if img_path and os.path.exists(img_path):
            pic = slide.shapes.add_picture(img_path, img_left, img_top, width=img_w, height=img_h)
            # 부드러운 가장자리 효과 (둥근 모서리 사각형으로 대체)
            pic.auto_shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
        else:
            # 실패 시 대체 상자
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, img_left, img_top, img_w, img_h)
            box.text = "이미지\n없음"

        # --- [텍스트 배치] ---
        # 이미지(1+7=8cm) 옆 1cm 띄우고 시작 = 9cm
        text_left = img_left + img_w + Cm(1) 
        text_width = slide_width - text_left - Cm(1) # 오른쪽 여백 1cm 남김
        text_height = Cm(6)
        text_top = Cm(1)

        txBox = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)
        tf = txBox.text_frame
        tf.word_wrap = True

        # 명언 내용
        p1 = tf.paragraphs[0]
        p1.text = quote
        p1.font.name = 'HY크리스탈M'
        p1.font.size = Pt(40)
        p1.alignment = PP_ALIGN.LEFT
        
        # 인물 이름
        p2 = tf.add_paragraph()
        p2.text = f"- {author}"
        p2.font.name = 'HY크리스탈M'
        p2.font.size = Pt(32)
        p2.alignment = PP_ALIGN.RIGHT
        p2.space_before = Pt(20)

    # 저장
    save_path = os.path.join(FOLDER_NAME, FILE_NAME)
    prs.save(save_path)
    print(f"\n[성공] 모든 작업이 완료되었습니다! 파일 위치: {save_path}")

if __name__ == "__main__":
    create_final_ppt()